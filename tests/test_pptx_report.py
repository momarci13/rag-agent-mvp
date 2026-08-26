"""Tests for tools/pptx_report.py: builds a well-formed deck with the
expected content, using python-pptx's own read API to inspect it."""

from pptx import Presentation

from agents.risk_schemas import ValidationFinding, ValidationReport
from tools.pptx_report import DISCLAIMER_BANNER, ValidationDeckBuilder


def _sample_report() -> ValidationReport:
    findings = [
        ValidationFinding(
            domain="credit_risk", area="Calibration and back-testing", verdict="non_compliant",
            severity="critical", description="PSI breach detected in the retail PD model case file.",
            recommendation="Recalibrate the model and reassess quarterly.", owner="Model Owner",
            remediation_deadline_days=30,
        ),
        ValidationFinding(
            domain="credit_risk", area="Conceptual soundness", verdict="compliant",
            severity="observation", description="Model design is well documented and industry-consistent.",
        ),
    ]
    return ValidationReport(
        domain="credit_risk", title="Credit Risk Validation Report -- PD-RETAIL-01",
        scope="Scope text", methodology="Methodology text",
        entity_under_review="PD-RETAIL-01 (mortgages)", reporting_period="2026Q2",
        findings=findings, quantitative_results={"psi": 0.30, "gini": 0.35},
        overall_rating="non_compliant", overall_conclusion="Critical PSI breach requires recalibration.",
    )


def _all_text(prs: Presentation) -> str:
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks)


def test_build_produces_a_well_formed_deck(tmp_path):
    report = _sample_report()
    output_path = ValidationDeckBuilder().build(report, tmp_path / "report.pptx")

    assert output_path.exists()
    assert output_path.stat().st_size > 1000

    prs = Presentation(str(output_path))
    assert len(prs.slides) >= 9  # title + 8 core sections at minimum

    full_text = _all_text(prs)
    assert DISCLAIMER_BANNER in full_text
    assert "PSI breach detected in the retail PD model case file" in full_text
    assert report.disclaimer in full_text
    assert "PENDING" in full_text  # not yet signed off


def test_build_reflects_signed_off_state(tmp_path):
    report = _sample_report()
    report.signed_off_by = "jane.validator"
    report.signed_off_at = "2026-06-30T00:00:00+00:00"
    output_path = ValidationDeckBuilder().build(report, tmp_path / "signed.pptx")
    prs = Presentation(str(output_path))
    full_text = _all_text(prs)
    assert "jane.validator" in full_text
    assert "PENDING" not in full_text
