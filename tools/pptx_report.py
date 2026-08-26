"""PPTX validation deck builder.

Builds a draft PowerPoint deck from a single, domain-agnostic
:class:`agents.risk_schemas.ValidationReport`. Every slide carries a visible
"DRAFT -- not a regulatory submission" banner, and the sign-off slide is
regenerated (not just re-labeled) once the report has been signed off --
see agents/risk_validation_team.py::RiskValidationOrchestrator.execute.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from agents.risk_schemas import ValidationFinding, ValidationReport

DISCLAIMER_BANNER = "DRAFT -- NOT A REGULATORY SUBMISSION"

_SEVERITY_COLOR = {
    "critical": RGBColor(0xC0, 0x00, 0x00),
    "high": RGBColor(0xE0, 0x6C, 0x0C),
    "medium": RGBColor(0xBF, 0x8F, 0x00),
    "low": RGBColor(0x54, 0x82, 0x35),
    "observation": RGBColor(0x59, 0x59, 0x59),
}

_MAX_FINDING_ROWS_PER_SLIDE = 8


class ValidationDeckBuilder:
    """Builds a draft PPTX validation deck from a ValidationReport."""

    def build(self, report: ValidationReport, output_path: str | Path) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._title_slide(prs, report)
        self._text_slide(prs, "Executive Summary", report.overall_conclusion or "(pending)")
        self._text_slide(
            prs, "Scope & Methodology",
            f"Scope:\n{report.scope}\n\nMethodology:\n{report.methodology}",
        )
        self._overview_slide(prs, report)
        self._activities_slide(prs, report)
        self._findings_slides(prs, report)
        self._quantitative_slide(prs, report)
        self._recommendations_slide(prs, report)
        self._conclusion_slide(prs, report)
        self._signoff_slide(prs, report)

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    # ---------- slide helpers ----------

    def _blank_slide(self, prs: Presentation, heading: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        self._banner(prs, slide)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), prs.slide_width - Inches(1.0), Inches(0.7))
        tf = title_box.text_frame
        tf.text = heading
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        return slide

    def _banner(self, prs: Presentation, slide) -> None:
        box = slide.shapes.add_textbox(Inches(0.3), prs.slide_height - Inches(0.4), prs.slide_width - Inches(0.6), Inches(0.3))
        tf = box.text_frame
        tf.text = DISCLAIMER_BANNER
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        run.font.bold = True

    def _title_slide(self, prs: Presentation, report: ValidationReport) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._banner(prs, slide)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), prs.slide_width - Inches(1.4), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = report.title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        meta_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), prs.slide_width - Inches(1.4), Inches(2.0))
        meta_tf = meta_box.text_frame
        meta_tf.word_wrap = True
        lines = [
            f"Entity under review: {report.entity_under_review}",
            f"Reporting period: {report.reporting_period}",
            f"Domain: {report.domain.replace('_', ' ').title()}",
            f"Generated: {report.generated_at}",
            f"Prepared by: {report.prepared_by}",
        ]
        meta_tf.text = lines[0]
        for line in lines[1:]:
            p = meta_tf.add_paragraph()
            p.text = line
        for p in meta_tf.paragraphs:
            p.font.size = Pt(14)

    def _text_slide(self, prs: Presentation, heading: str, body: str) -> None:
        slide = self._blank_slide(prs, heading)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), prs.slide_width - Inches(1.0), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = body
        for p in tf.paragraphs:
            p.font.size = Pt(16)

    def _overview_slide(self, prs: Presentation, report: ValidationReport) -> None:
        body = (
            f"Entity under review: {report.entity_under_review}\n"
            f"Reporting period: {report.reporting_period}\n"
            f"Domain: {report.domain.replace('_', ' ').title()}\n\n"
            f"Case file summary (as submitted):\n"
        )
        for key, value in report.quantitative_results.items():
            body += f"  - {key}: {value}\n"
        self._text_slide(prs, "Model / Exposure Overview", body)

    def _activities_slide(self, prs: Presentation, report: ValidationReport) -> None:
        slide = self._blank_slide(prs, "Validation Activities Checklist")
        areas: dict[str, str] = {}
        for f in report.findings:
            if f.area not in areas or f.verdict == "non_compliant":
                areas[f.area] = f.verdict
        rows = max(1, len(areas)) + 1
        table_shape = slide.shapes.add_table(rows, 2, Inches(0.5), Inches(1.3), Inches(12.0), Inches(0.5) * rows)
        table = table_shape.table
        table.columns[0].width = Inches(9.0)
        table.columns[1].width = Inches(3.0)
        table.cell(0, 0).text = "Area"
        table.cell(0, 1).text = "Verdict"
        if not areas:
            table.cell(1, 0).text = "No validation areas assessed"
            table.cell(1, 1).text = "n/a"
        else:
            for i, (area, verdict) in enumerate(areas.items(), start=1):
                table.cell(i, 0).text = area
                table.cell(i, 1).text = verdict

    def _findings_slides(self, prs: Presentation, report: ValidationReport) -> None:
        findings = report.findings
        if not findings:
            self._text_slide(prs, "Findings & Severity Ratings", "No findings recorded.")
            return
        for start in range(0, len(findings), _MAX_FINDING_ROWS_PER_SLIDE):
            chunk = findings[start:start + _MAX_FINDING_ROWS_PER_SLIDE]
            heading = "Findings & Severity Ratings"
            if len(findings) > _MAX_FINDING_ROWS_PER_SLIDE:
                heading += f" ({start // _MAX_FINDING_ROWS_PER_SLIDE + 1})"
            self._findings_table_slide(prs, heading, chunk)

    def _findings_table_slide(self, prs: Presentation, heading: str, findings: list[ValidationFinding]) -> None:
        slide = self._blank_slide(prs, heading)
        rows = len(findings) + 1
        table_shape = slide.shapes.add_table(rows, 4, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.5) * rows)
        table = table_shape.table
        widths = [Inches(2.6), Inches(1.6), Inches(1.2), Inches(7.1)]
        for i, w in enumerate(widths):
            table.columns[i].width = w
        for i, header in enumerate(["Area", "Verdict", "Severity", "Description"]):
            table.cell(0, i).text = header
        for row, f in enumerate(findings, start=1):
            table.cell(row, 0).text = f.area
            table.cell(row, 1).text = f.verdict
            cell = table.cell(row, 2)
            cell.text = f.severity
            color = _SEVERITY_COLOR.get(f.severity)
            if color is not None and cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = color
            table.cell(row, 3).text = f.description[:280]

    def _quantitative_slide(self, prs: Presentation, report: ValidationReport) -> None:
        slide = self._blank_slide(prs, "Quantitative Test Results")
        results = report.quantitative_results
        if not results:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11.0), Inches(1.0))
            box.text_frame.text = "No quantitative results were supplied."
            return
        rows = len(results) + 1
        table_shape = slide.shapes.add_table(rows, 2, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.4) * rows)
        table = table_shape.table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        for i, (key, value) in enumerate(results.items(), start=1):
            table.cell(i, 0).text = str(key)
            table.cell(i, 1).text = str(value)

    def _recommendations_slide(self, prs: Presentation, report: ValidationReport) -> None:
        actionable = [f for f in report.findings if f.recommendation]
        slide = self._blank_slide(prs, "Recommendations & Remediation Plan")
        if not actionable:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11.0), Inches(1.0))
            box.text_frame.text = "No outstanding recommendations."
            return
        rows = len(actionable) + 1
        table_shape = slide.shapes.add_table(rows, 3, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.5) * rows)
        table = table_shape.table
        widths = [Inches(7.0), Inches(2.75), Inches(2.75)]
        for i, w in enumerate(widths):
            table.columns[i].width = w
        for i, header in enumerate(["Recommendation", "Owner", "Deadline (days)"]):
            table.cell(0, i).text = header
        for row, f in enumerate(actionable, start=1):
            table.cell(row, 0).text = f.recommendation[:280]
            table.cell(row, 1).text = f.owner or "Unassigned"
            table.cell(row, 2).text = str(f.remediation_deadline_days) if f.remediation_deadline_days is not None else "n/a"

    def _conclusion_slide(self, prs: Presentation, report: ValidationReport) -> None:
        slide = self._blank_slide(prs, "Overall Validation Conclusion")
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), prs.slide_width - Inches(1.4), Inches(1.0))
        tf = box.text_frame
        tf.text = f"Overall rating: {report.overall_rating.upper()}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), prs.slide_width - Inches(1.4), Inches(3.5))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.text = report.overall_conclusion or "(pending)"
        for p in body_tf.paragraphs:
            p.font.size = Pt(16)

    def _signoff_slide(self, prs: Presentation, report: ValidationReport) -> None:
        slide = self._blank_slide(prs, "Sign-off")
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), prs.slide_width - Inches(1.4), Inches(1.6))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"Prepared by: {report.prepared_by}"
        p2 = tf.add_paragraph()
        p2.text = f"Signed off by: {report.signed_off_by or 'PENDING -- human validator sign-off required'}"
        p3 = tf.add_paragraph()
        p3.text = f"Signed off at: {report.signed_off_at or 'n/a'}"
        for p in tf.paragraphs:
            p.font.size = Pt(16)

        disclaimer_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), prs.slide_width - Inches(1.4), Inches(3.2))
        d_tf = disclaimer_box.text_frame
        d_tf.word_wrap = True
        d_tf.text = report.disclaimer
        d_tf.paragraphs[0].font.size = Pt(13)
        d_tf.paragraphs[0].font.italic = True
